"""Convert uploaded content files (PDF, PPTX, DOCX) into safe HTML for in-browser preview."""

import html as _html
import base64
import importlib.util
import os

PREVIEW_DEPENDENCIES = {
    'pptx': ('pptx', 'python-pptx'),
    'docx': ('docx', 'python-docx'),
}


def get_preview_dependency_status() -> dict[str, bool]:
    status = {}
    for key, (module_name, _package_name) in PREVIEW_DEPENDENCIES.items():
        status[key] = importlib.util.find_spec(module_name) is not None
    return status


def get_missing_preview_dependencies() -> list[str]:
    missing = []
    for key, installed in get_preview_dependency_status().items():
        if not installed:
            missing.append(PREVIEW_DEPENDENCIES[key][1])
    return missing


def preview_exception_message(ext: str, exc: Exception) -> str:
    if isinstance(exc, ModuleNotFoundError):
        key = 'docx' if ext in ('doc', 'docx') else ext
        package_name = PREVIEW_DEPENDENCIES.get(key, (None, None))[1]
        if package_name:
            return (
                f'{ext.upper()} preview is unavailable because the server is missing '
                f'the optional dependency `{package_name}`.'
            )
    return str(exc)


# ── PPTX → HTML ─────────────────────────────────────────────────────────────

def pptx_to_html(file_path: str) -> str:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    slides = []

    for slide_num, slide in enumerate(prs.slides, 1):
        parts = []

        for shape in slide.shapes:
            # ── Text frames ──
            if shape.has_text_frame:
                tf = shape.text_frame
                is_title = (
                    hasattr(shape, 'placeholder_format')
                    and shape.placeholder_format is not None
                    and shape.placeholder_format.idx == 0
                )
                for para in tf.paragraphs:
                    if not para.text.strip():
                        continue
                    runs_html = []
                    for run in para.runs:
                        t = _html.escape(run.text)
                        if run.font.bold:
                            t = f'<strong>{t}</strong>'
                        if run.font.italic:
                            t = f'<em>{t}</em>'
                        if run.font.size:
                            px = min(int(run.font.size.pt), 36)
                            t = f'<span style="font-size:{px}px">{t}</span>'
                        runs_html.append(t)
                    inner = ''.join(runs_html) or _html.escape(para.text)
                    if is_title:
                        parts.append(f'<h4 class="slide-title">{inner}</h4>')
                    else:
                        parts.append(f'<p class="slide-para">{inner}</p>')

            # ── Embedded images ──
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                    ct   = shape.image.content_type or 'image/png'
                    b64  = base64.b64encode(blob).decode()
                    parts.append(
                        f'<img src="data:{ct};base64,{b64}" '
                        f'class="slide-img" alt="slide image">'
                    )
                except Exception:
                    pass

        body = '\n'.join(parts) if parts else '<p class="text-muted small">Empty slide</p>'
        slides.append(
            f'<div class="pptx-slide" id="slide-{slide_num}">'
            f'<div class="slide-num-badge">Slide {slide_num} / {len(prs.slides)}</div>'
            f'<div class="slide-content">{body}</div>'
            f'</div>'
        )

    return '\n'.join(slides)


# ── DOCX → HTML ──────────────────────────────────────────────────────────────

def docx_to_html(file_path: str) -> str:
    import docx as _docx
    from docx.oxml.ns import qn

    doc    = _docx.Document(file_path)
    parts  = []
    in_ul  = False
    in_ol  = False

    def close_list():
        nonlocal in_ul, in_ol
        if in_ul:
            parts.append('</ul>')
            in_ul = False
        if in_ol:
            parts.append('</ol>')
            in_ol = False

    for para in doc.paragraphs:
        style = para.style.name
        text  = para.text

        if not text.strip():
            close_list()
            continue

        # Build inline-formatted content from runs
        runs = []
        for run in para.runs:
            t = _html.escape(run.text)
            if run.bold:
                t = f'<strong>{t}</strong>'
            if run.italic:
                t = f'<em>{t}</em>'
            if run.underline:
                t = f'<u>{t}</u>'
            runs.append(t)
        inner = ''.join(runs) if runs else _html.escape(text)

        if 'Heading 1' in style:
            close_list()
            parts.append(f'<h2 class="docx-h1">{inner}</h2>')
        elif 'Heading 2' in style:
            close_list()
            parts.append(f'<h3 class="docx-h2">{inner}</h3>')
        elif 'Heading 3' in style or 'Heading 4' in style:
            close_list()
            parts.append(f'<h4 class="docx-h3">{inner}</h4>')
        elif 'List Bullet' in style:
            if not in_ul:
                close_list()
                parts.append('<ul class="docx-list">')
                in_ul = True
            parts.append(f'<li>{inner}</li>')
        elif 'List Number' in style:
            if not in_ol:
                close_list()
                parts.append('<ol class="docx-list">')
                in_ol = True
            parts.append(f'<li>{inner}</li>')
        else:
            close_list()
            parts.append(f'<p class="docx-para">{inner}</p>')

    close_list()

    # Tables (appended after paragraphs — order within the doc is not preserved
    # for tables mixed with paragraphs, but works for most simple documents)
    for table in doc.tables:
        parts.append(
            '<div class="table-responsive my-3">'
            '<table class="table table-bordered table-sm align-middle">'
        )
        for i, row in enumerate(table.rows):
            parts.append('<tr>')
            tag = 'th' if i == 0 else 'td'
            for cell in row.cells:
                parts.append(f'<{tag}>{_html.escape(cell.text)}</{tag}>')
            parts.append('</tr>')
        parts.append('</table></div>')

    return '\n'.join(parts)
